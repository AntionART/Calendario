from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

presentation = Presentation()
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
title_tf = title_box.text_frame
title_tf.text = 'Cómo usar la aplicación Calendario'
title_tf.paragraphs[0].font.size = Pt(36)
title_tf.paragraphs[0].font.bold = True

subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(9), Inches(0.6))
subtitle_tf = subtitle_box.text_frame
subtitle_tf.text = 'Guía rápida para Usuarios y Administradores'
subtitle_tf.paragraphs[0].font.size = Pt(24)
subtitle_tf.paragraphs[0].font.italic = True
subtitle_tf.paragraphs[0].font.color.rgb = RGBColor(11, 120, 209)

users_box = slide.shapes.add_textbox(Inches(0.4), Inches(1.9), Inches(4.4), Inches(5.4))
users_tf = users_box.text_frame
users_tf.word_wrap = True
users_tf.text = 'Usuarios'
users_tf.paragraphs[0].font.size = Pt(28)
users_tf.paragraphs[0].font.bold = True
users_tf.paragraphs[0].font.color.rgb = RGBColor(0, 104, 183)

users_points = [
    'Acceso: Abra la app en el navegador e inicie sesión con su usuario y contraseña.',
    'Ver agenda: Seleccione “Calendario” para ver turnos por día/semana/mes.',
    'Reservar turno: Click en “Nuevo turno” → elegir fecha, hora, tipo y confirmar.',
    'Modificar/Cancelar: Abrir turno reservado → “Editar” o “Cancelar”.',
    'Notificaciones: Revise la campana para mensajes y recordatorios.',
    'Buenas prácticas: Verifique disponibilidad y confirme datos del paciente.'
]
for point in users_points:
    p = users_tf.add_paragraph()
    p.text = point
    p.level = 1
    p.font.size = Pt(18)
    p.space_after = Pt(6)

admin_box = slide.shapes.add_textbox(Inches(5), Inches(1.9), Inches(4.4), Inches(5.4))
admin_tf = admin_box.text_frame
admin_tf.word_wrap = True
admin_tf.text = 'Administradores'
admin_tf.paragraphs[0].font.size = Pt(28)
admin_tf.paragraphs[0].font.bold = True
admin_tf.paragraphs[0].font.color.rgb = RGBColor(0, 160, 157)

admin_points = [
    'Panel admin: Iniciar sesión con rol administrador y acceder a “Administración”.',
    'Gestión de usuarios: Crear/editar roles (médico, enfermero, recepción).',
    'Configurar horarios: Definir franjas, turnos por profesional y días hábiles.',
    'Crear/Asignar turnos manualmente: Desde “Turnos” → “Agregar” → asignar profesional.',
    'Sincronizar contador: Usar después de restauraciones o reinicios.',
    'Exportar/reportes: Generar listados diarios y exportar a CSV/PDF.'
]
for point in admin_points:
    p = admin_tf.add_paragraph()
    p.text = point
    p.level = 1
    p.font.size = Pt(18)
    p.space_after = Pt(6)

notes = slide.notes_slide.notes_text_frame
notes.text = (
    'Objetivo: explicar el uso principal del Calendario para usuarios y administradores.\n'
    'Usuarios:\n'
    '- Mostrar reservas, edición y cancelación de turnos.\n'
    '- Recordar las notificaciones y verificar datos.\n'
    'Administradores:\n'
    '- Explicar cómo gestionar usuarios, horarios y turnos.\n'
    '- Mencionar sincronización del contador y exportación de reportes.\n'
    'Cierre: invitar a usar la app con cuidado y contactar soporte si hay dudas.'
)

presentation.save('Diapositiva_Calendario.pptx')
