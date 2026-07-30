from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

presentation = Presentation()
presentation.slide_width = Inches(10)
presentation.slide_height = Inches(7.5)

slide_layout = presentation.slide_layouts[6]  # Blank layout
slide = presentation.slides.add_slide(slide_layout)

# Fondo blanco
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 255, 255)

# Colores
COLOR_BLUE = RGBColor(11, 120, 209)
COLOR_GREEN = RGBColor(0, 160, 157)
COLOR_DARK = RGBColor(40, 40, 40)

# === TITULO ===
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.8))
title_tf = title_box.text_frame
title_tf.text = 'Cómo usar la Aplicación Calendario'
title_tf.paragraphs[0].font.size = Pt(44)
title_tf.paragraphs[0].font.bold = True
title_tf.paragraphs[0].font.color.rgb = COLOR_DARK

# Línea divisoria
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.1), Inches(9), Inches(0.03))
line.fill.solid()
line.fill.fore_color.rgb = COLOR_BLUE
line.line.color.rgb = COLOR_BLUE

# === USUARIOS SECTION (LEFT) ===
user_box_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.35), Inches(1.3), Inches(4.7), Inches(2.2))
user_box_shape.fill.solid()
user_box_shape.fill.fore_color.rgb = RGBColor(230, 240, 250)
user_box_shape.line.color.rgb = COLOR_BLUE
user_box_shape.line.width = Pt(2)

# Icono usuario
user_icon = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), Inches(1.5), Inches(0.5), Inches(0.5))
user_icon.fill.solid()
user_icon.fill.fore_color.rgb = COLOR_BLUE
user_icon.line.width = Pt(0)
user_text = user_icon.text_frame
user_text.text = "👤"
user_text.paragraphs[0].font.size = Pt(24)

# Título usuarios
user_title = slide.shapes.add_textbox(Inches(1.2), Inches(1.55), Inches(3.8), Inches(0.4))
user_title_tf = user_title.text_frame
user_title_tf.text = 'USUARIOS'
user_title_tf.paragraphs[0].font.size = Pt(22)
user_title_tf.paragraphs[0].font.bold = True
user_title_tf.paragraphs[0].font.color.rgb = COLOR_BLUE

# Puntos usuarios
user_text_box = slide.shapes.add_textbox(Inches(0.55), Inches(2.05), Inches(4.4), Inches(1.35))
user_text_frame = user_text_box.text_frame
user_text_frame.word_wrap = True

user_points = [
    '✓ Acceso inmediato: Abre la app en el navegador',
    '✓ Ver agenda: Consulta turnos por día/semana/mes',
    '✓ Reservar: "Nuevo turno" → fecha, hora, tipo',
    '✓ Editar/Cancelar: Modifica tu turno cuando sea necesario',
    '✓ Notificaciones: Recibe recordatorios y mensajes'
]

for i, point in enumerate(user_points):
    if i == 0:
        user_text_frame.text = point
        user_text_frame.paragraphs[0].font.size = Pt(15)
        user_text_frame.paragraphs[0].font.color.rgb = COLOR_DARK
    else:
        p = user_text_frame.add_paragraph()
        p.text = point
        p.font.size = Pt(15)
        p.font.color.rgb = COLOR_DARK
        p.space_before = Pt(1)
        p.space_after = Pt(1)

# === ADMINISTRADORES SECTION (RIGHT) ===
admin_box_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.15), Inches(1.3), Inches(4.5), Inches(2.2))
admin_box_shape.fill.solid()
admin_box_shape.fill.fore_color.rgb = RGBColor(230, 250, 248)
admin_box_shape.line.color.rgb = COLOR_GREEN
admin_box_shape.line.width = Pt(2)

# Icono administrador
admin_icon = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.35), Inches(1.5), Inches(0.5), Inches(0.5))
admin_icon.fill.solid()
admin_icon.fill.fore_color.rgb = COLOR_GREEN
admin_icon.line.width = Pt(0)
admin_text = admin_icon.text_frame
admin_text.text = "⚙️"
admin_text.paragraphs[0].font.size = Pt(24)

# Título administradores
admin_title = slide.shapes.add_textbox(Inches(5.95), Inches(1.55), Inches(3.6), Inches(0.4))
admin_title_tf = admin_title.text_frame
admin_title_tf.text = 'ADMINISTRADORES'
admin_title_tf.paragraphs[0].font.size = Pt(22)
admin_title_tf.paragraphs[0].font.bold = True
admin_title_tf.paragraphs[0].font.color.rgb = COLOR_GREEN

# Puntos administradores
admin_text_box = slide.shapes.add_textbox(Inches(5.3), Inches(2.05), Inches(4.2), Inches(1.35))
admin_text_frame = admin_text_box.text_frame
admin_text_frame.word_wrap = True

admin_points = [
    '✓ Panel admin: Gestiona usuarios y roles',
    '✓ Horarios: Configura franjas y profesionales',
    '✓ Crear turnos: Asigna turnos manualmente',
    '✓ Sincronizador: Verifica contador tras reinicios',
    '✓ Reportes: Exporta datos a CSV/PDF'
]

for i, point in enumerate(admin_points):
    if i == 0:
        admin_text_frame.text = point
        admin_text_frame.paragraphs[0].font.size = Pt(15)
        admin_text_frame.paragraphs[0].font.color.rgb = COLOR_DARK
    else:
        p = admin_text_frame.add_paragraph()
        p.text = point
        p.font.size = Pt(15)
        p.font.color.rgb = COLOR_DARK
        p.space_before = Pt(1)
        p.space_after = Pt(1)

# === CARACTERÍSTICAS PRINCIPALES (BOTTOM) ===
features_title = slide.shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(9), Inches(0.35))
features_title_tf = features_title.text_frame
features_title_tf.text = 'Características Principales'
features_title_tf.paragraphs[0].font.size = Pt(20)
features_title_tf.paragraphs[0].font.bold = True
features_title_tf.paragraphs[0].font.color.rgb = COLOR_DARK

# Feature boxes
features = [
    {'icon': '📅', 'title': 'Calendario Inteligente', 'desc': 'Visualiza turnos por día, semana o mes'},
    {'icon': '🔔', 'title': 'Notificaciones', 'desc': 'Recordatorios automáticos de citas'},
    {'icon': '💾', 'title': 'Datos Seguros', 'desc': 'Información sincronizada en tiempo real'},
    {'icon': '📊', 'title': 'Reportes', 'desc': 'Exporta estadísticas y listados'}
]

feature_x_start = 0.6
feature_width = 2.15

for idx, feature in enumerate(features):
    x_pos = feature_x_start + (idx * feature_width)
    
    # Feature box
    feat_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_pos), Inches(4.15), Inches(2.0), Inches(1.4))
    feat_box.fill.solid()
    feat_box.fill.fore_color.rgb = RGBColor(245, 245, 245)
    feat_box.line.color.rgb = RGBColor(200, 200, 200)
    feat_box.line.width = Pt(1)
    
    # Icon
    icon_box = slide.shapes.add_textbox(Inches(x_pos + 0.3), Inches(4.25), Inches(1.4), Inches(0.4))
    icon_box_tf = icon_box.text_frame
    icon_box_tf.text = feature['icon']
    icon_box_tf.paragraphs[0].font.size = Pt(28)
    
    # Title
    title_box_feat = slide.shapes.add_textbox(Inches(x_pos + 0.15), Inches(4.75), Inches(1.7), Inches(0.4))
    title_box_feat_tf = title_box_feat.text_frame
    title_box_feat_tf.text = feature['title']
    title_box_feat_tf.paragraphs[0].font.size = Pt(12)
    title_box_feat_tf.paragraphs[0].font.bold = True
    title_box_feat_tf.paragraphs[0].font.color.rgb = COLOR_DARK
    title_box_feat_tf.word_wrap = True
    
    # Description
    desc_box = slide.shapes.add_textbox(Inches(x_pos + 0.15), Inches(5.2), Inches(1.7), Inches(0.3))
    desc_box_tf = desc_box.text_frame
    desc_box_tf.text = feature['desc']
    desc_box_tf.paragraphs[0].font.size = Pt(9)
    desc_box_tf.paragraphs[0].font.color.rgb = RGBColor(80, 80, 80)
    desc_box_tf.word_wrap = True

# === FOOTER ===
footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.9), Inches(9), Inches(0.5))
footer_tf = footer_box.text_frame
footer_tf.text = '💡 Consejo: La app se ejecuta localmente en tu navegador. ¡No requiere configuración adicional! | 📧 Soporte: Contacta al administrador si tienes dudas'
footer_tf.paragraphs[0].font.size = Pt(11)
footer_tf.paragraphs[0].font.italic = True
footer_tf.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)
footer_tf.word_wrap = True

# === NOTAS DEL ORADOR ===
notes = slide.notes_slide.notes_text_frame
notes.text = (
    'GUÍA DE PRESENTACIÓN\n\n'
    'Introducción (10 segundos):\n'
    '- La aplicación Calendario te permite gestionar turnos de forma fácil y rápida.\n'
    '- Se ejecuta localmente en tu navegador. ¡Acceso inmediato sin configuración!\n\n'
    'Usuarios (30 segundos):\n'
    '- Pueden ver la agenda completa, reservar nuevos turnos, modificarlos o cancelarlos.\n'
    '- Reciben notificaciones automáticas para no olvidar sus citas.\n'
    '- Interfaz simple y amigable.\n\n'
    'Administradores (30 segundos):\n'
    '- Panel completo para gestionar usuarios, roles y horarios.\n'
    '- Crear turnos manualmente cuando sea necesario.\n'
    '- Sincronización automática y exportación de reportes.\n'
    '- Herramientas de control y auditoría.\n\n'
    'Características (20 segundos):\n'
    '- Calendario inteligente con múltiples vistas.\n'
    '- Notificaciones en tiempo real.\n'
    '- Datos seguros y sincronizados.\n'
    '- Reportes y estadísticas.\n\n'
    'Cierre (10 segundos):\n'
    '- La app está lista para usar sin configuración ni credenciales.\n'
    '- Invita preguntas y ofrece contacto de soporte.'
)

# Guardar con nuevo nombre
presentation.save('Diapositiva_Calendario_v2.pptx')
print("✅ Diapositiva mejorada creada: Diapositiva_Calendario_v2.pptx")
print("   ✓ Acceso sin usuario/contraseña")
print("   ✓ Iconos y emojis visuales")
print("   ✓ Diseño con 2 columnas + características")
print("   ✓ Footer con consejo importante")
